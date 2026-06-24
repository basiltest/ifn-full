#!/usr/bin/env python3
"""
Internship Project Report I deck for Basil Stevenson (Lumenor).
Same institutional / "issued, not branded" look as build_deck.py, but the slides
follow the report-format section structure (cover + report sections, lean).
Image placeholders are positioned for screenshots Basil drops into img/.
No em-dashes or en-dashes anywhere (repo rule): hyphens / words only.
"""
import os, glob
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(HERE, "img")
BUILD_DIR = os.path.join(IMG_DIR, ".build")
os.makedirs(BUILD_DIR, exist_ok=True)

TOTAL = 11   # total slides, for footer page numbers

# ----- palette -------------------------------------------------------------
INK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0xB1, 0x12, 0x26)   # ICFAI red (matches IFN logo bar)
PANEL  = RGBColor(0xF4, 0xF2, 0xEE)   # warm off-white
CHIP   = RGBColor(0xF7, 0xF5, 0xF1)
HAIR   = RGBColor(0xDE, 0xDA, 0xD2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Carlito"

SW, SH = Inches(13.333), Inches(7.5)
MX     = Inches(0.95)          # left text margin
CW     = Inches(11.55)         # content width
BAND   = Inches(0.16)          # left accent band width

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----- low-level helpers ---------------------------------------------------
def _flatten(shp):
    """Remove the themed <p:style> (kills effectRef drop-shadow) and disable shadow,
    so shapes render flat in both PowerPoint and LibreOffice."""
    el = shp._element
    st = el.find(qn('p:style'))
    if st is not None:
        el.remove(st)
    shp.shadow.inherit = False


def slide():
    s = prs.slides.add_slide(BLANK)
    # white background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    _flatten(bg)
    # left accent band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, BAND, SH)
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT
    band.line.fill.background(); _flatten(band)
    return s


def _set_dash(line, val="dash"):
    ln = line._get_or_add_ln()
    d = ln.find(qn('a:prstDash'))
    if d is None:
        d = ln.makeelement(qn('a:prstDash'), {})
        ln.append(d)
    d.set('val', val)


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, dash=None, rounded=False, radius=0.06):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
        if dash:
            _set_dash(shp.line, dash)
    _flatten(shp)
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=4,
         letter=None):
    """runs: str, or list of paragraphs; each paragraph is str or list of (txt, opts) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        segs = para if isinstance(para, list) else [(para, {})]
        for txt, opts in segs:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic)
            f.color.rgb = opts.get("color", color)
            if letter is not None or "letter" in opts:
                _spc(r, opts.get("letter", letter))
    return tb


def _spc(run, pts):
    if pts is None: return
    run._r.get_or_add_rPr().set('spc', str(int(pts * 100)))


def header(s, kicker, title):
    # kicker
    text(s, MX, Inches(0.62), CW, Inches(0.3),
         kicker.upper(), size=12.5, color=ACCENT, bold=True, letter=1.4)
    # title
    text(s, MX, Inches(0.95), CW, Inches(0.9),
         title, size=30, color=INK, bold=True)
    # hairline
    rect(s, MX, Inches(1.78), CW, Pt(1.4), fill=HAIR)
    return Inches(2.05)


def footer(s, n):
    rect(s, MX, Inches(6.96), CW, Pt(1.0), fill=HAIR)
    text(s, MX, Inches(7.04), Inches(8), Inches(0.3),
         "Basil Stevenson  ·  Lumenor  ·  Internship Review",
         size=9, color=MUTED)
    text(s, Inches(9.9), Inches(7.04), Inches(2.49), Inches(0.3),
         "%02d / %02d" % (n, TOTAL), size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=14, gap=8, lead_color=INK, body_color=INK):
    """items: list of (lead, body). lead bold; body normal. lead may be ''."""
    paras = []
    for lead, body in items:
        seg = [("•  ", {"color": ACCENT, "bold": True, "size": size})]
        if lead:
            seg.append((lead, {"color": lead_color, "bold": True, "size": size}))
        seg.append((body, {"color": body_color, "size": size}))
        paras.append(seg)
    return text(s, x, y, w, h, paras, size=size, spacing=1.04, space_after=gap)


def placeholder(s, x, y, w, h, title, sub):
    box = rect(s, x, y, w, h, fill=PANEL, line=HAIR, lw=1.4, dash="dash", rounded=True, radius=0.04)
    text(s, x, y + h/2 - Inches(0.42), w, Inches(0.4),
         "▢  " + title, size=13, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.3), y + h/2 + Inches(0.02), w - Inches(0.6), Inches(0.5),
         sub, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    return box


def _find_img(key):
    for f in sorted(glob.glob(os.path.join(IMG_DIR, key + ".*"))):
        if os.path.isfile(f) and not f.endswith((".build",)):
            return f
    return None


def slot(s, x, y, w, h, key, ph_title, ph_sub, anchor="center"):
    """Embed img/<key>.* cover-cropped to the box (no distortion); else dashed placeholder.
    anchor controls which part survives a crop: 'top'/'bottom' (vertical), 'left'/'right' (horizontal)."""
    src = _find_img(key)
    if not src:
        return placeholder(s, x, y, w, h, ph_title, ph_sub)
    im = Image.open(src).convert("RGB")
    iw, ih = im.size
    target = w / h                      # box aspect (EMU / EMU)
    if iw / ih > target:                # image too wide -> crop sides
        nw = int(round(ih * target))
        x0 = 0 if anchor == "left" else (iw - nw) if anchor == "right" else (iw - nw) // 2
        im = im.crop((x0, 0, x0 + nw, ih))
    else:                               # image too tall -> crop top/bottom
        nh = int(round(iw / target))
        y0 = 0 if anchor == "top" else (ih - nh) if anchor == "bottom" else (ih - nh) // 2
        im = im.crop((0, y0, iw, y0 + nh))
    out = os.path.join(BUILD_DIR, key + ".png")
    im.save(out, "PNG")
    pic = s.shapes.add_picture(out, x, y, width=w, height=h)
    pic.line.color.rgb = HAIR; pic.line.width = Pt(1.0)
    return pic


def stat_tile(s, x, y, w, h, number, label):
    rect(s, x, y, w, h, fill=PANEL, rounded=True, radius=0.08)
    rect(s, x, y, Inches(0.09), h, fill=ACCENT)  # accent edge
    text(s, x + Inches(0.28), y + Inches(0.16), w - Inches(0.4), Inches(0.7),
         number, size=33, color=INK, bold=True)
    text(s, x + Inches(0.3), y + h - Inches(0.56), w - Inches(0.4), Inches(0.5),
         label, size=11.5, color=MUTED)


def chip(s, x, y, label, size=12.5):
    w = Inches(0.42 + 0.092 * len(label))
    rect(s, x, y, w, Inches(0.46), fill=CHIP, line=HAIR, lw=1.0, rounded=True, radius=0.5)
    text(s, x, y + Inches(0.085), w, Inches(0.3), label, size=size, color=INK,
         align=PP_ALIGN.CENTER)
    return w


def chip_grid(s, x, y, labels, row_max=Inches(12.4), row_h=Inches(0.6)):
    cx = x; cy = y
    for lab in labels:
        w = Inches(0.42 + 0.095 * len(lab))
        if cx + w > row_max:
            cx = x; cy = cy + row_h
        chip(s, cx, cy, lab)
        cx = cx + Inches(0.42 + 0.095 * len(lab)) + Inches(0.18)
    return cy + row_h


# ===========================================================================
# SLIDE 1 - Cover page
# ===========================================================================
s = slide()
text(s, MX, Inches(1.2), CW, Inches(0.4),
     "INTERNSHIP PROJECT REPORT I", size=14, color=ACCENT, bold=True, letter=2.0)
text(s, MX, Inches(1.7), CW, Inches(1.4),
     "Lumenor Internship Review", size=52, color=INK, bold=True)
text(s, MX, Inches(2.85), CW, Inches(0.5),
     "Lumenor   ·   May to July 2026", size=21, color=MUTED)
rect(s, MX, Inches(3.6), Inches(2.0), Pt(2.2), fill=ACCENT)
text(s, MX, Inches(3.85), Inches(10.8), Inches(0.8),
     "Three work-streams: QA on FlowX, a product landing page, and the IFN platform.",
     size=17, color=INK)
# identity block (cover-page fields)
rect(s, MX, Inches(5.35), CW, Pt(1.2), fill=HAIR)
text(s, MX, Inches(5.55), Inches(8), Inches(0.4),
     "Basil Stevenson", size=17, color=INK, bold=True)
text(s, MX, Inches(5.98), Inches(11), Inches(1.0),
     [[("Roll No 24BCAHH010052   ·   BCA, Second Year   ·   ICFAI", {"color": MUTED, "size": 13})],
      [("Company mentor: Rama Raju sir", {"color": MUTED, "size": 13})],
      [("Internship duration: May to July 2026", {"color": MUTED, "size": 13})]],
     space_after=2)

# ===========================================================================
# SLIDE 2 - Section 1: Introduction
# ===========================================================================
s = slide()
y = header(s, "Section 1", "Introduction")
text(s, MX, y, CW, Inches(0.9),
     [[("Lumenor", {"bold": True}),
       (" builds software for solar-energy businesses: a customer suite (", {}),
       ("solarX, flowX, fixX", {"bold": True}),
       ("), plus ", {}),
       ("coreX", {"bold": True}),
       (" for internal teams.", {})]],
     size=15, color=INK, spacing=1.06)
bullets(s, MX, Inches(2.9), CW, Inches(3.2), [
    ("Team assigned:  ", "the FlowX team: Inventory, Sales and CRM, HRMS, Accounting, Procurement."),
    ("Project domain:  ", "Quality Assurance and full-stack web development."),
    ("Purpose of the internship:  ", "test real product flows, ship production frontend, and build a platform end to end."),
], size=15, gap=14)
footer(s, 2)

# ===========================================================================
# SLIDE 3 - Section 2: Project overview (problem statement folded in)
# ===========================================================================
s = slide()
y = header(s, "Section 2", "Project overview")
cards = [
    ("01", "QA on FlowX", "Flow-test FlowX and flag flaws. Problem: ship a reliable product before release."),
    ("02", "Lumenor landing page", "Build the public marketing site. Problem: present the product to prospects."),
    ("03", "IFN platform", "Build a members-only founder platform. Problem: connect the university founder community."),
]
cw = Inches(3.62); gap = Inches(0.30); cx = MX; cy = Inches(2.3); ch = Inches(2.7)
for num, t, d in cards:
    rect(s, cx, cy, cw, ch, fill=PANEL, rounded=True, radius=0.05)
    rect(s, cx, cy, cw, Inches(0.10), fill=ACCENT)
    text(s, cx + Inches(0.3), cy + Inches(0.32), cw - Inches(0.6), Inches(0.5),
         num, size=15, color=ACCENT, bold=True)
    text(s, cx + Inches(0.3), cy + Inches(0.74), cw - Inches(0.6), Inches(0.6),
         t, size=18, color=INK, bold=True)
    text(s, cx + Inches(0.3), cy + Inches(1.4), cw - Inches(0.6), Inches(1.2),
         d, size=13, color=MUTED, spacing=1.05)
    cx = cx + cw + gap
rect(s, MX, Inches(5.45), Inches(0.09), Inches(0.5), fill=ACCENT)
text(s, MX + Inches(0.25), Inches(5.47), CW, Inches(0.5),
     [[("Objective: ", {"bold": True, "color": INK}),
       ("contribute across QA, frontend, and full-stack product within one internship.", {"color": INK})]],
     size=14)
footer(s, 3)

# ===========================================================================
# SLIDE 4 - Section 5: Work completed - QA on FlowX
# ===========================================================================
s = slide()
y = header(s, "Section 5  ·  Work completed", "QA on FlowX")
tiles = [("67", "tests run"), ("19", "bugs found"),
         ("7", "modules covered"), ("2", "security criticals")]
tw = Inches(2.74); tgap = Inches(0.18); tx = MX; ty = Inches(2.12); th = Inches(1.25)
for num, lab in tiles:
    stat_tile(s, tx, ty, tw, th, num, lab)
    tx = tx + tw + tgap
text(s, MX, Inches(3.5), CW, Inches(0.6),
     [[("Brief:  ", {"bold": True, "color": INK}),
       ("test FlowX like a user, optimize the happy paths, flag flaws. Assigned to FixX first; moved to FlowX when FixX was not ready. I also offered security testing, from a penetration-testing background.",
        {"color": INK})]],
     size=12.5, spacing=1.05)
text(s, MX, Inches(4.32), Inches(6.6), Inches(0.35),
     "Selected findings (mine)", size=14, color=ACCENT, bold=True)
bullets(s, MX, Inches(4.7), Inches(6.55), Inches(1.9), [
    ("Account takeover:  ", "password-reset poisoning via a client-controlled baseUrl leaks the reset token to an attacker domain (Critical)."),
    ("Info disclosure:  ", "a duplicate warehouse code surfaces the raw ORM error, leaking internal file paths and schema."),
    ("Data integrity and crashes:  ", "address corrupted on every save, a broken CSV round-trip, a search that crashes on an undefined field."),
], size=12.5, gap=6)
slot(s, Inches(7.85), Inches(4.32), Inches(4.55), Inches(2.02),
     "bug-tracker", "Bug tracker", "Screenshot: shared BUGS sheet + Playwright report", anchor="top")
text(s, MX, Inches(6.5), CW, Inches(0.4),
     [[("Tooling:  ", {"bold": True, "color": INK}),
       ("Playwright automated flows (shared HTML reports); each bug logged with module, repro steps, expected versus actual, priority, status.",
        {"color": MUTED})]],
     size=11.5)
footer(s, 4)

# ===========================================================================
# SLIDE 5 - Section 5: Work completed - Lumenor landing page
# ===========================================================================
s = slide()
y = header(s, "Section 5  ·  Work completed", "Lumenor landing page")
bullets(s, MX, Inches(2.2), Inches(5.5), Inches(3.5), [
    ("", "Built and shipped the public marketing site for Lumenor."),
    ("", "Headline: “Run your entire solar business in one place.”"),
    ("", "Responsive layout that presents the product to prospects."),
], size=15, gap=14)
text(s, MX, Inches(4.55), Inches(5.5), Inches(0.4),
     [[("Live:  ", {"bold": True, "color": INK}),
       ("basiltest.github.io/lumenor", {"color": ACCENT, "bold": True})]],
     size=14)
slot(s, Inches(6.7), Inches(2.2), Inches(5.7), Inches(3.95),
     "landing", "Landing page", "Screenshot: the landing page hero")
footer(s, 5)

# ===========================================================================
# SLIDE 6 - Section 5: Work completed - IFN platform
# ===========================================================================
s = slide()
y = header(s, "Section 5  ·  Work completed", "IFN: ICFAI Founders Network")
text(s, MX, Inches(2.0), CW, Inches(0.5),
     "A members-only platform for the university's founder community, built and iterated largely solo.",
     size=15, color=INK)
labels = ["Feed", "Idea Pipeline G1-G6", "Problem Hub", "Team Acquisition", "Calendar",
          "Directory", "Notifications", "Moderation", "Admin", "Auth and Roles"]
chip_grid(s, MX, Inches(2.5), labels)
rect(s, MX, Inches(3.62), CW, Inches(0.85), fill=PANEL, rounded=True, radius=0.05)
rect(s, MX, Inches(3.62), Inches(0.09), Inches(0.85), fill=ACCENT)
text(s, MX + Inches(0.3), Inches(3.74), CW - Inches(0.6), Inches(0.7),
     [[("Stack:  ", {"bold": True, "color": INK}),
       ("Vite single-page app on Vercel, backed by Supabase (Auth + Postgres + Row-Level Security), no custom backend. Now hardening for production and moving to self-hosted infrastructure.",
        {"color": INK})]],
     size=13.5, spacing=1.05)
slot(s, MX, Inches(4.68), Inches(5.62), Inches(2.16),
     "ifn-feed", "IFN feed", "Screenshot: feed with posts, voting, tags", anchor="top")
slot(s, Inches(6.83), Inches(4.68), Inches(5.62), Inches(2.16),
     "ifn-pipeline", "Idea pipeline", "Screenshot: pipeline dossier / gate bar", anchor="top")
footer(s, 6)

# ===========================================================================
# SLIDE 7 - Section 6: Methodology and approach
# ===========================================================================
s = slide()
y = header(s, "Section 6", "Methodology and approach")
bullets(s, MX, y, CW, Inches(4.0), [
    ("Process:  ", "Agile and iterative work, with weekly progress and review inside the team."),
    ("QA approach:  ", "exploratory testing plus Playwright automation; every bug logged with module, repro steps, expected versus actual, priority, and status."),
    ("Build approach:  ", "full-stack delivery: auth, Row-Level Security, data modeling, and a single-page app shipped to production."),
    ("Working practice:  ", "dogfooding the products inside a live SaaS team."),
], size=15, gap=16)
footer(s, 7)

# ===========================================================================
# SLIDE 8 - Section 7: Tools and technologies
# ===========================================================================
s = slide()
y = header(s, "Section 7", "Tools and technologies")
text(s, MX, y, CW, Inches(0.5),
     "Languages, frameworks, data, and tooling used across the three work-streams.",
     size=14, color=MUTED)
tools = ["Python", "JavaScript", "React", "Vite", "Playwright", "Supabase", "Postgres",
         "Row-Level Security", "Vercel", "Docker", "Git and GitHub", "LLMs", "Vector databases"]
chip_grid(s, MX, Inches(2.85), tools)
footer(s, 8)

# ===========================================================================
# SLIDE 9 - Section 8: Results and progress status
# ===========================================================================
s = slide()
y = header(s, "Section 8", "Results and progress status")
tiles = [("19", "bugs logged"), ("2", "security criticals"),
         ("3", "products delivered"), ("~55%", "internship progress")]
tw = Inches(2.74); tgap = Inches(0.18); tx = MX; ty = Inches(2.12); th = Inches(1.3)
for num, lab in tiles:
    stat_tile(s, tx, ty, tw, th, num, lab)
    tx = tx + tw + tgap
bullets(s, MX, Inches(3.85), CW, Inches(2.6), [
    ("QA:  ", "67 tests run across 7 FlowX modules; 19 bugs logged, including 2 security criticals reported with repro steps."),
    ("Landing page:  ", "shipped live to prospects."),
    ("IFN platform:  ", "built end to end and iterating; auth, RLS, feed, idea pipeline, and admin in place."),
], size=14.5, gap=12)
rect(s, MX, Inches(6.0), Inches(0.09), Inches(0.5), fill=ACCENT)
text(s, MX + Inches(0.25), Inches(6.02), CW, Inches(0.5),
     [[("Project completion status: ", {"bold": True, "color": INK}),
       ("about 55 percent, at the mid-internship point.", {"color": INK})]],
     size=14)
footer(s, 9)

# ===========================================================================
# SLIDE 10 - Section 9: Challenges and solutions
# ===========================================================================
s = slide()
y = header(s, "Section 9", "Challenges and solutions")
bullets(s, MX, y, CW, Inches(4.2), [
    ("FixX not ready:  ", "reassigned to FlowX and ramped on its modules quickly to keep testing on schedule."),
    ("Surfacing security flaws:  ", "reported auth and error-handling bugs responsibly, each with clear repro steps and priority."),
    ("Solo scope on IFN:  ", "sliced the build into vertical features and shipped iteratively instead of all at once."),
    ("Self-host migration:  ", "planned a phased move off Supabase Cloud and Vercel onto the office server, with captcha to follow."),
], size=15, gap=16)
footer(s, 10)

# ===========================================================================
# SLIDE 11 - Sections 10 to 12: Outcomes, work plan, conclusion
# ===========================================================================
s = slide()
y = header(s, "Sections 10 to 12", "Outcomes, plan, and conclusion")
colw = Inches(5.55)
text(s, MX, Inches(2.15), colw, Inches(0.35),
     "Work plan (second half)", size=14, color=ACCENT, bold=True)
bullets(s, MX, Inches(2.55), colw, Inches(3.0), [
    ("IFN:  ", "finish the pipeline end-to-end testing; moderation and reporting."),
    ("IFN:  ", "self-host migration off Supabase Cloud and Vercel onto the office server; captcha rollout."),
    ("FlowX:  ", "retest the fixed bugs and widen coverage."),
], size=14, gap=12)
rx = Inches(6.95)
text(s, rx, Inches(2.15), colw, Inches(0.35),
     "Learning outcomes", size=14, color=ACCENT, bold=True)
bullets(s, rx, Inches(2.55), colw, Inches(3.0), [
    ("", "Reproducible bug reporting, priority triage, and security testing of auth and error handling."),
    ("", "Full-stack delivery: auth, RLS and data modeling, a SPA, shipped to production."),
    ("", "Working and dogfooding inside a real SaaS team."),
    ("", "Building depth in Docker, Git, agentic AI, LLMs, and vector databases."),
], size=14, gap=11)
rect(s, MX, Inches(5.95), Inches(0.09), Inches(0.5), fill=ACCENT)
text(s, MX + Inches(0.25), Inches(5.97), CW, Inches(0.5),
     "Conclusion: halfway in, bugs found, products shipped, more to build.",
     size=15, color=INK, bold=True)
footer(s, 11)

# ---------------------------------------------------------------------------
out = "/home/basil/lumenor/ifn/presentation/Lumenor_Internship_Report.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
